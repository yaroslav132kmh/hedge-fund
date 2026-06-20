"""Build PowerPoint and PDF presentations (RU + EN) with diagrams from solution artifacts."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "_presentation_assets"
OUT = ROOT / "docs"

# ── colours (quant / dark-blue palette) ──────────────────────────────────────
C_TITLE = RGBColor(0x15, 0x65, 0xC0)
C_BODY = RGBColor(0x26, 0x32, 0x38)
C_MUTED = RGBColor(0x78, 0x90, 0x9C)
C_ACCENT = RGBColor(0x4F, 0xC3, 0xF7)

AGENT_STEPS = [
    ("DataAcquisition", "Data"),
    ("MarketAnalysis", "Analysis"),
    ("HestonCalibration", "Heston"),
    ("GreeksCalculation", "Greeks"),
    ("HedgingDecision", "Hedge"),
    ("PortfolioOptimization", "Portfolio"),
    ("RiskManagement", "Risk"),
    ("Backtesting", "Backtest"),
    ("SelfDiagnostic", "Diagnostic"),
    ("Explainability", "Explain"),
    ("Dashboard", "Dashboard"),
]

STRATEGY_STEPS_RU = [
    ("1. Liability", "Short put +\nvega call\n(опционная книга)"),
    ("2. Size hedge", "Risk budget\n2% daily VaR\n→ BTC notional"),
    ("3. Delta hedge", "Spot BTC\nposition\n(Δ → 0)"),
    ("4. Vega hedge", "ATM call\n(ν → 0)"),
    ("5. Portfolio", "15 assets\nRisk Parity\nrebal. 5d"),
    ("6. Risk layer", "VaR/CVaR\nstops\ntrailing"),
]

STRATEGY_STEPS_EN = [
    ("1. Liability", "Short put +\nvega call\n(option book)"),
    ("2. Size hedge", "Risk budget\n2% daily VaR\n→ BTC notional"),
    ("3. Delta hedge", "Spot BTC\nposition\n(Δ → 0)"),
    ("4. Vega hedge", "ATM call\n(ν → 0)"),
    ("5. Portfolio", "15 assets\nRisk Parity\nrebal. 5d"),
    ("6. Risk layer", "VaR/CVaR\nstops\ntrailing"),
]


def _load_context_data():
    """Load artifacts from results/ or run pipeline from checkpoints."""
    results = ROOT / "artifacts" / "results"
    data = {}
    paths = {
        "hedge_history": results / "hedging_history.parquet",
        "portfolio_equity": results / "portfolio_equity.parquet",
        "portfolio_weights": results / "portfolio_weights_path.parquet",
        "portfolio_methods": results / "portfolio_methods.parquet",
        "stress": results / "stress_test.parquet",
        "spot_close": ROOT / "artifacts" / "checkpoints" / "default",
    }
    if paths["hedge_history"].exists():
        data["hedge_history"] = pd.read_parquet(paths["hedge_history"])
    if paths["portfolio_equity"].exists():
        data["portfolio_equity"] = pd.read_parquet(paths["portfolio_equity"])
    if paths["portfolio_weights"].exists():
        data["portfolio_weights"] = pd.read_parquet(paths["portfolio_weights"])
    if paths["portfolio_methods"].exists():
        data["portfolio_methods"] = pd.read_parquet(paths["portfolio_methods"])
    if paths["stress"].exists():
        data["stress"] = pd.read_parquet(paths["stress"])

    # spot from checkpoint or bundled data
    spot_ckpt = ROOT / "artifacts" / "checkpoints" / "default" / "data_acquisition__spot_close.parquet"
    if spot_ckpt.exists():
        data["spot_close"] = pd.read_parquet(spot_ckpt)
    else:
        bundled = ROOT / "data" / "raw" / "spot_close.parquet"
        if bundled.exists():
            data["spot_close"] = pd.read_parquet(bundled)

    calibr = ROOT / "artifacts" / "calibration" / "calibr_data.parquet"
    if calibr.exists():
        data["calibr"] = pd.read_parquet(calibr)

    if "hedge_history" not in data or data["hedge_history"].empty:
        print("Running pipeline to generate artifacts...", file=sys.stderr)
        from cryptohedge.agents import build_pipeline
        from cryptohedge.core.config import load_config
        from cryptohedge.core.context import AgentContext
        cfg = load_config(ROOT / "config")
        ctx = AgentContext(cfg, root=ROOT)
        build_pipeline(ctx, fail_fast=True).run()
        return _load_context_data()
    return data


def _style_ax(ax, title=""):
    ax.set_facecolor("#1a1a2e")
    ax.figure.patch.set_facecolor("#1a1a2e")
    ax.tick_params(colors="#cfd8dc", labelsize=8)
    ax.title.set_color("#4fc3f7")
    ax.title.set_fontsize(11)
    ax.set_title(title)
    for spine in ax.spines.values():
        spine.set_color("#37474f")
    ax.grid(True, alpha=0.2, color="#546e7a")
    ax.xaxis.label.set_color("#90a4ae")
    ax.yaxis.label.set_color("#90a4ae")


def draw_agents_flow(path: Path, lang: str):
    fig, ax = plt.subplots(figsize=(14, 3.2))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 3)
    ax.axis("off")
    fig.patch.set_facecolor("#0f1115")
    n = len(AGENT_STEPS)
    xs = np.linspace(0.4, 13.2, n)
    for i, (name, short) in enumerate(AGENT_STEPS):
        x, y = xs[i], 1.5
        box = FancyBboxPatch((x - 0.52, y - 0.42), 1.04, 0.84,
                             boxstyle="round,pad=0.04", facecolor="#1565c0",
                             edgecolor="#4fc3f7", linewidth=1.2)
        ax.add_patch(box)
        ax.text(x, y + 0.08, short, ha="center", va="center", fontsize=8,
                color="white", fontweight="bold")
        ax.text(x, y - 0.22, name[:12], ha="center", va="center", fontsize=5.5,
                color="#b0bec5")
        if i < n - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.55, y), xytext=(x + 0.55, y),
                        arrowprops=dict(arrowstyle="->", color="#66bb6a", lw=1.5))
    title = "Конвейер AI-агентов CryptoHedge" if lang == "ru" else "CryptoHedge AI Agent Pipeline"
    ax.text(7, 2.65, title, ha="center", fontsize=13, color="#4fc3f7", fontweight="bold")
    ax.text(7, 0.35, "Blackboard + MessageBus  |  Orchestrator  |  Checkpoint/resume",
            ha="center", fontsize=9, color="#78909c")
    fig.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="#0f1115")
    plt.close(fig)


def draw_strategy(path: Path, lang: str):
    steps = STRATEGY_STEPS_RU if lang == "ru" else STRATEGY_STEPS_EN
    fig, ax = plt.subplots(figsize=(12, 3.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 3.5)
    ax.axis("off")
    fig.patch.set_facecolor("#0f1115")
    n = len(steps)
    xs = np.linspace(0.8, 11.2, n)
    colours = ["#c62828", "#ef6c00", "#1565c0", "#00838f", "#2e7d32", "#6a1b9a"]
    for i, (hdr, body) in enumerate(steps):
        x, y = xs[i], 1.6
        box = FancyBboxPatch((x - 0.75, y - 0.7), 1.5, 1.4,
                             boxstyle="round,pad=0.05", facecolor=colours[i],
                             edgecolor="white", linewidth=0.8, alpha=0.9)
        ax.add_patch(box)
        ax.text(x, y + 0.35, hdr, ha="center", va="center", fontsize=8,
                color="white", fontweight="bold")
        ax.text(x, y - 0.2, body, ha="center", va="center", fontsize=7, color="white")
        if i < n - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.8, y), xytext=(x + 0.8, y),
                        arrowprops=dict(arrowstyle="->", color="#ffca28", lw=2))
    t = "Стратегия: хеджирование + портфель" if lang == "ru" else "Strategy: hedging + portfolio"
    ax.text(6, 3.1, t, ha="center", fontsize=13, color="#4fc3f7", fontweight="bold")
    sub = ("Heston-модель · Δ-ν нейтрализация · Risk Parity · адаптивные стопы"
           if lang == "ru" else
           "Heston model · Δ-ν neutralisation · Risk Parity · adaptive stops")
    ax.text(6, 0.25, sub, ha="center", fontsize=9, color="#78909c")
    fig.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="#0f1115")
    plt.close(fig)


def draw_solution_charts(data: dict, lang: str):
    charts = {}

    # spot
    if "spot_close" in data:
        sc = data["spot_close"]
        primary = "BTCUSDT" if "BTCUSDT" in sc.columns else sc.columns[0]
        fig, ax = plt.subplots(figsize=(10, 3))
        idx = pd.to_datetime(sc.index)
        ax.plot(idx, sc[primary], color="#4fc3f7", lw=1.5)
        _style_ax(ax, "BTCUSDT" if lang == "en" else "Цена BTC (спот)")
        fig.tight_layout()
        p = ASSETS / f"spot_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor=ax.get_facecolor())
        plt.close(fig)
        charts["spot"] = p

    # hedge PnL
    if "hedge_history" in data:
        h = data["hedge_history"]
        ts = pd.to_datetime(h["ts"])
        fig, axes = plt.subplots(2, 1, figsize=(10, 5), sharex=True)
        axes[0].plot(ts, h["spot"], color="#4fc3f7")
        axes[0].set_ylabel("USD")
        _style_ax(axes[0], "Spot" if lang == "en" else "Спот")
        axes[1].plot(ts, h["pnl"], color="#66bb6a", label="PnL")
        axes[1].plot(ts, h["fee"], color="#ffa726", label="Fees" if lang == "en" else "Комиссии")
        axes[1].legend(facecolor="#263238", labelcolor="white", fontsize=8)
        _style_ax(axes[1], "Hedge PnL" if lang == "en" else "PnL хедж-стратегии")
        fig.tight_layout()
        p = ASSETS / f"hedge_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
        plt.close(fig)
        charts["hedge"] = p

    # portfolio equity + weights
    if "portfolio_equity" in data:
        eq = data["portfolio_equity"]
        ts = pd.to_datetime(eq["ts"])
        fig, ax = plt.subplots(figsize=(10, 3.5))
        ax.plot(ts, eq["equity"], color="#66bb6a", lw=2,
                label="Portfolio" if lang == "en" else "Портфель")
        if "benchmark" in eq.columns:
            ax.plot(ts, eq["benchmark"], color="#90a4ae", ls="--",
                    label="EQ benchmark" if lang == "en" else "Бенчмарк")
        ax.legend(facecolor="#263238", labelcolor="white", fontsize=8)
        _style_ax(ax, "Portfolio equity" if lang == "en" else "Стоимость портфеля")
        fig.tight_layout()
        p = ASSETS / f"equity_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
        plt.close(fig)
        charts["equity"] = p

    if "portfolio_weights" in data:
        wp = data["portfolio_weights"]
        ts = pd.to_datetime(wp["ts"])
        cols = [c for c in wp.columns if c != "ts"]
        means = wp[cols].mean().sort_values(ascending=False)
        show = list(means.head(10).index)
        fig, ax = plt.subplots(figsize=(10, 3.5))
        ax.stackplot(ts, *[wp[c] for c in show], labels=show, alpha=0.85)
        ax.set_ylim(0, 1)
        ax.legend(loc="upper left", fontsize=6, ncol=2, facecolor="#263238", labelcolor="white")
        _style_ax(ax, "Weight evolution" if lang == "en" else "Эволюция весов")
        fig.tight_layout()
        p = ASSETS / f"weights_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
        plt.close(fig)
        charts["weights"] = p

    # method comparison
    if "portfolio_methods" in data:
        m = data["portfolio_methods"]
        fig, ax1 = plt.subplots(figsize=(10, 3.5))
        colours = ["#66bb6a" if c else "#455a64" for c in m.get("chosen", [False] * len(m))]
        ax1.bar(m["method"], m["total_return"], color=colours)
        ax1.set_ylabel("Return" if lang == "en" else "Доходность")
        ax1.tick_params(axis="x", rotation=25, labelsize=7)
        _style_ax(ax1, "Optimization methods" if lang == "en" else "Методы оптимизации")
        fig.tight_layout()
        p = ASSETS / f"methods_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
        plt.close(fig)
        charts["methods"] = p

    # stress test
    if "stress" in data:
        s = data["stress"]
        if "net_hedged_pnl" in s.columns:
            fig, ax = plt.subplots(figsize=(10, 3.5))
            x = np.arange(len(s))
            w = 0.35
            ax.bar(x - w / 2, s["unhedged_pnl"], w, color="#ef5350",
                   label="Unhedged" if lang == "en" else "Без хеджа")
            ax.bar(x + w / 2, s["net_hedged_pnl"], w, color="#66bb6a",
                   label="Hedged" if lang == "en" else "С хеджем")
            ax.set_xticks(x)
            ax.set_xticklabels(s["scenario"], rotation=20, fontsize=7)
            ax.legend(facecolor="#263238", labelcolor="white", fontsize=8)
            _style_ax(ax, "Stress tests" if lang == "en" else "Стресс-тесты")
            fig.tight_layout()
            p = ASSETS / f"stress_{lang}.png"
            fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
            plt.close(fig)
            charts["stress"] = p

    # heston params
    if "calibr" in data:
        cal = data["calibr"].sort_values("sample_idx")
        t = cal["sample_idx"]
        fig, axes = plt.subplots(2, 3, figsize=(10, 4))
        for ax, p in zip(axes.flat, ["v0", "kappa", "theta", "eps", "rho"]):
            ax.plot(t, cal[p], color="#ab47bc", lw=1.2)
            _style_ax(ax, p)
        axes.flat[-1].axis("off")
        fig.suptitle("Heston parameters" if lang == "en" else "Параметры Хестона",
                     color="#4fc3f7", fontsize=11)
        fig.tight_layout()
        p = ASSETS / f"heston_{lang}.png"
        fig.savefig(p, dpi=120, bbox_inches="tight", facecolor="#1a1a2e")
        plt.close(fig)
        charts["heston"] = p

    return charts


# ── slide content ─────────────────────────────────────────────────────────────
SLIDES_RU = [
    {"title": "CryptoHedge", "subtitle": "Мультиагентный криптовалютный хедж-фонд · MVP",
     "bullets": ["11 автономных AI-агентов", "Капитал $10M · горизонт 90 дней",
                 "Heston-хедж BTC + портфель 15 активов", "42 теста · бэктест · дашборды RU/EN"]},
    {"title": "Торговая вселенная и данные",
     "bullets": ["100 спот USDT-пар; опционы BTC 30d (call/put, 11 страйков)",
                 "Провайдеры: bundled · synthetic · binance (live REST)",
                 "Комиссии: spot 0.03%; опционы Deribit-style 0.03% + cap 12.5%",
                 "Масштабирование: MarketDataProvider → OKX/Bybit/Deribit"]},
    {"title": "Мультиагентная архитектура",
     "bullets": ["Blackboard (AgentContext) — без прямых вызовов между агентами",
                 "MessageBus — типизированные сообщения, correlation_id",
                 "Clean Architecture: agents → services → domain → pyquant",
                 "Checkpoint/resume; конфиг YAML + Pydantic"]},
    {"title": "Схема агентов", "image": "agents", "bullets": []},
    {"title": "Стратегия хеджирования и портфеля", "image": "strategy",
     "bullets": ["Short put (0.95 K/S) + vega-call — опционная книга (liability)",
                 "Размер хеджа из risk budget 2% дневного VaR",
                 "Δ-хедж спотом, ν-хедж опционом → остаточные греки ≈ 0",
                 "Портфель: Risk Parity, 15 активов, ребаланс 5d, max weight 20%"]},
    {"title": "Количественный движок",
     "bullets": ["Heston: MLE + IV surface; контроль устойчивости; бенчмарки BS/SABR",
                 "Греки: Δ Γ ν Θ ρ vanna volga charm",
                 "Стресс −10% BTC: unhedged ≈ −$108k, hedged ≈ $0"],
     "image": "heston"},
    {"title": "Управление рисками",
     "bullets": ["VaR (hist/gauss/CF), CVaR, MDD limit 25%, leverage 3×",
                 "Защищено: Δ, ν, концентрация, tail shocks",
                 "Не защищено: execution, counterparty, basis, jumps",
                 "Допущения: daily bars, historical VaR, long-only, model fees"]},
    {"title": "Защитные механизмы",
     "bullets": ["Адаптивный стоп: ATR + VaR + Heston-vol (2–25%)",
                 "Трейлинг-стоп ATR×2.5", "Зоны Δ: green 5% / red 15%",
                 "Confidence Score: drift + calibration + hedge quality"]},
    {"title": "Управление портфелем",
     "bullets": ["5 оптимизаторов; авто-выбор Risk Parity",
                 "MVP: return +42.1%, DR 1.63, effective N 13.7/15",
                 "Бэктест с drift весов и комиссиями"],
     "image": "equity", "image2": "weights", "image3": "methods"},
    {"title": "Результаты хеджа и стресс-тесты",
     "bullets": ["Walk-forward: train 30d / test 5d; bias controls",
                 "Остаточная Δ → 0 на каждом срезе",
                 "seed=90909090 — воспроизводимость всех RNG"],
     "image": "hedge", "image2": "stress"},
    {"title": "Валидация и roadmap",
     "bullets": ["42 автотеста; solution.ipynb end-to-end; max|ΔPnL|=0",
                 "Roadmap: live execution multi-CEX, Telegram, sentiment, RL",
                 "Межбиржевой арбитраж · live опционы · meta-coordination"],
     "image": "spot"},
]

SLIDES_EN = [
    {"title": "CryptoHedge", "subtitle": "Multi-agent cryptocurrency hedge fund · MVP",
     "bullets": ["11 autonomous AI agents", "Capital $10M · 90-day horizon",
                 "Heston BTC hedge + 15-asset portfolio", "42 tests · backtest · RU/EN dashboards"]},
    {"title": "Trading universe & data",
     "bullets": ["100 spot USDT pairs; BTC options 30d (call/put, 11 strikes)",
                 "Providers: bundled · synthetic · binance (live REST)",
                 "Fees: spot 0.03%; Deribit-style options 0.03% + cap 12.5%",
                 "Scaling: MarketDataProvider → OKX/Bybit/Deribit"]},
    {"title": "Multi-agent architecture",
     "bullets": ["Blackboard (AgentContext) — no direct agent-to-agent calls",
                 "MessageBus — typed messages, correlation_id",
                 "Clean Architecture: agents → services → domain → pyquant",
                 "Checkpoint/resume; YAML config + Pydantic"]},
    {"title": "Agent pipeline", "image": "agents", "bullets": []},
    {"title": "Hedging & portfolio strategy", "image": "strategy",
     "bullets": ["Short put (0.95 K/S) + vega call — option liability book",
                 "Hedge size from 2% daily VaR risk budget",
                 "Δ hedge via spot, ν hedge via option → residual greeks ≈ 0",
                 "Portfolio: Risk Parity, 15 assets, 5d rebalance, max weight 20%"]},
    {"title": "Quantitative engine",
     "bullets": ["Heston: MLE + IV surface; stability control; BS/SABR benchmarks",
                 "Greeks: Δ Γ ν Θ ρ vanna volga charm",
                 "Stress −10% BTC: unhedged ≈ −$108k, hedged ≈ $0"],
     "image": "heston"},
    {"title": "Risk management",
     "bullets": ["VaR (hist/gauss/CF), CVaR, MDD limit 25%, leverage 3×",
                 "Protected: Δ, ν, concentration, tail shocks",
                 "Not protected: execution, counterparty, basis, jumps",
                 "Assumptions: daily bars, historical VaR, long-only, model fees"]},
    {"title": "Protective mechanisms",
     "bullets": ["Adaptive stop: ATR + VaR + Heston-vol (2–25%)",
                 "Trailing stop ATR×2.5", "Δ zones: green 5% / red 15%",
                 "Confidence Score: drift + calibration + hedge quality"]},
    {"title": "Portfolio management",
     "bullets": ["5 optimizers; auto-selected Risk Parity",
                 "MVP: return +42.1%, DR 1.63, effective N 13.7/15",
                 "Backtest with weight drift and fees"],
     "image": "equity", "image2": "weights", "image3": "methods"},
    {"title": "Hedge results & stress tests",
     "bullets": ["Walk-forward: train 30d / test 5d; bias controls",
                 "Residual Δ → 0 each slice",
                 "seed=90909090 — reproducibility of all RNG"],
     "image": "hedge", "image2": "stress"},
    {"title": "Validation & roadmap",
     "bullets": ["42 automated tests; solution.ipynb end-to-end; max|ΔPnL|=0",
                 "Roadmap: live multi-CEX execution, Telegram, sentiment, RL",
                 "Cross-exchange arb · live options · meta-coordination"],
     "image": "spot"},
]


def _add_title_slide(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = C_MUTED


def _add_bullets(slide, bullets, top=1.3, left=0.6, width=12, size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(size)
        p.font.color.rgb = C_BODY
        p.space_after = Pt(6)


def _resolve_image(key: str, lang: str, charts: dict) -> Path | None:
    if key == "agents":
        return ASSETS / f"agents_flow_{lang}.png"
    if key == "strategy":
        return ASSETS / f"strategy_{lang}.png"
    return charts.get(key)


def build_pptx(slides: list, charts: dict, lang: str, out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for spec in slides:
        slide = prs.slides.add_slide(blank)
        _add_title_slide(slide, spec["title"], spec.get("subtitle", ""))

        imgs = []
        for k in ("image", "image2", "image3"):
            if spec.get(k):
                p = _resolve_image(spec[k], lang, charts)
                if p and p.exists():
                    imgs.append(p)

        if imgs and not spec.get("bullets"):
            # full-width diagram slide
            slide.shapes.add_picture(str(imgs[0]), Inches(0.4), Inches(1.2), width=Inches(12.5))
        elif imgs:
            top = 1.15
            if spec.get("bullets"):
                _add_bullets(slide, spec["bullets"], top=top, width=5.8, size=12)
            # images on the right or bottom
            if len(imgs) == 1:
                slide.shapes.add_picture(str(imgs[0]), Inches(6.5), Inches(1.2), width=Inches(6.3))
            elif len(imgs) == 2:
                slide.shapes.add_picture(str(imgs[0]), Inches(6.5), Inches(1.2), width=Inches(6.3))
                slide.shapes.add_picture(str(imgs[1]), Inches(6.5), Inches(4.0), width=Inches(6.3))
            else:
                slide.shapes.add_picture(str(imgs[0]), Inches(6.5), Inches(1.2), width=Inches(6.3))
                slide.shapes.add_picture(str(imgs[1]), Inches(0.5), Inches(4.5), width=Inches(6.0))
                slide.shapes.add_picture(str(imgs[2]), Inches(6.8), Inches(4.5), width=Inches(6.0))
        elif spec.get("bullets"):
            _add_bullets(slide, spec["bullets"], top=1.3)

    prs.save(str(out_path))
    print(f"  PPTX -> {out_path}")


def build_pdf(slides: list, charts: dict, lang: str, out_path: Path):
    page_w, page_h = landscape(A4)
    doc = SimpleDocTemplate(str(out_path), pagesize=landscape(A4),
                            leftMargin=0.5 * inch, rightMargin=0.5 * inch,
                            topMargin=0.4 * inch, bottomMargin=0.4 * inch)
    styles = getSampleStyleSheet()
    title_st = ParagraphStyle("T", parent=styles["Heading1"], fontSize=22,
                              textColor=colors.HexColor("#1565C0"), spaceAfter=8)
    bullet_st = ParagraphStyle("B", parent=styles["Normal"], fontSize=11,
                               textColor=colors.HexColor("#263238"), leading=14,
                               leftIndent=12, bulletIndent=0)
    story = []

    for spec in slides:
        story.append(Paragraph(spec["title"], title_st))
        if spec.get("subtitle"):
            story.append(Paragraph(f'<font color="#78909c">{spec["subtitle"]}</font>',
                                   styles["Normal"]))
            story.append(Spacer(1, 6))
        for b in spec.get("bullets", []):
            story.append(Paragraph(f"• {b}", bullet_st))
        for k in ("image", "image2", "image3"):
            if spec.get(k):
                p = _resolve_image(spec[k], lang, charts)
                if p and p.exists():
                    story.append(Spacer(1, 8))
                    story.append(Image(str(p), width=9 * inch, height=3.2 * inch))
        story.append(PageBreak())

    doc.build(story)
    print(f"  PDF  -> {out_path}")


def try_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    """Try PowerPoint COM export on Windows."""
    try:
        import comtypes.client  # type: ignore
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        pres = powerpoint.Presentations.Open(str(pptx_path.resolve()))
        pres.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF
        pres.Close()
        powerpoint.Quit()
        print(f"  PDF (PowerPoint COM) -> {pdf_path}")
        return True
    except Exception as e:
        print(f"  PowerPoint COM export skipped: {e}", file=sys.stderr)
        return False


def main():
    ASSETS.mkdir(parents=True, exist_ok=True)
    print("Loading artifacts...")
    data = _load_context_data()

    for lang in ("ru", "en"):
        print(f"\n=== {lang.upper()} ===")
        draw_agents_flow(ASSETS / f"agents_flow_{lang}.png", lang)
        draw_strategy(ASSETS / f"strategy_{lang}.png", lang)
        charts = draw_solution_charts(data, lang)

        slides = SLIDES_RU if lang == "ru" else SLIDES_EN
        pptx_path = OUT / f"presentation.{lang}.pptx"
        pdf_path = OUT / f"presentation.{lang}.pdf"

        build_pptx(slides, charts, lang, pptx_path)
        if not try_pptx_to_pdf(pptx_path, pdf_path):
            build_pdf(slides, charts, lang, pdf_path)

    print("\nDone.")


if __name__ == "__main__":
    main()
